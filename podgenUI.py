"""
文本转语音生成器 (PodgenUI)

功能介绍:
该程序提供了一个基于 Streamlit 的用户界面，允许用户从文件或 URL 提取文本内容，并使用大模型处理文本。处理后的文本通过 TTS 技术生成音频文件，并支持音频预览和下载。

主要功能:
1. 支持从多种文件格式（txt, docx, pdf, html, md, pptx）提取文本
2. 支持从网页 URL 提取文本内容
3. 集成 GPT 模型进行文本处理（可选）
4. 使用 Microsoft Edge TTS 生成高质量语音
5. 提供多种语音选择
6. 支持音频预览和下载

使用方法:
1. 确保已安装所有依赖库，包括 streamlit、edge_tts、pydub、requests、html2text、pypdf、pptx、bs4、aiohttp 等。
   可以使用以下命令安装:
   pip install streamlit edge_tts pydub requests html2text pypdf pptx bs4 aiohttp

2. 运行程序:
   streamlit run podgenUI.py

3. 在浏览器中打开程序界面:
   - 上传文件或输入 URL 以获取文本内容
   - 使用 GPT 处理文本（可选）
   - 选择语音并生成音频
   - 预览和下载生成的音频文件

注意事项:
- 配置文件是podgenUI.ini
- 使用 GPT 功能需要配置 OpenAI API Key
- 确保网络连接正常，特别是从 URL 获取内容时
- 生成音频可能需要一些时间，请耐心等待

依赖:
- Streamlit
- Microsoft Edge TTS
- OpenAI API (可选)
- 其他 Python 库如 pydub、requests 等
"""

import os
import edge_tts
from pydub import AudioSegment
import streamlit as st
import asyncio
import tempfile
import docx
import PyPDF2
import openai
from io import BytesIO
import configparser
import json
import requests
from bs4 import BeautifulSoup
import re
from pptx import Presentation

def load_config():
    """从配置文件加载GPT设置"""
    config = configparser.ConfigParser()
    if os.path.exists('podgenUI.ini'):
        config.read('podgenUI.ini')
        if 'GPT' in config:
            return {
                'use_gpt': config['GPT'].getboolean('use_gpt', False),
                'model_name': config['GPT'].get('model_name', "gpt-3.5-turbo"),
                'base_url': config['GPT'].get('base_url', "https://api.openai.com/v1"),
                'api_key': config['GPT'].get('api_key', ""),
                'system_prompt': config['GPT'].get('system_prompt', "请将输入的文本整理成标题和正文的格式，去除不必要的内容，保持文章的主要信息。以纯文本格式输出。")
            }
    return {
        'use_gpt': False,
        'model_name': "gpt-3.5-turbo",
        'base_url': "https://api.openai.com/v1",
        'api_key': "",
        'system_prompt': "请将输入的文本整理成标题和正文的格式，去除不必要的内容，保持文章的主要信息。以纯文本格式输出。"
    }

def save_config(settings):
    """保存GPT设置到配置文件"""
    config = configparser.ConfigParser()
    config['GPT'] = {
        'use_gpt': str(settings['use_gpt']),
        'model_name': settings['model_name'],
        'base_url': settings['base_url'],
        'api_key': settings['api_key'],
        'system_prompt': settings['system_prompt']
    }
    with open('podgenUI.ini', 'w') as configfile:
        config.write(configfile)

# 加载配置
config_settings = load_config()

# 初始化session state
if 'openai_api_key' not in st.session_state:
    st.session_state.openai_api_key = config_settings['api_key']
if 'use_gpt' not in st.session_state:
    st.session_state.use_gpt = config_settings['use_gpt']
if 'model_name' not in st.session_state:
    st.session_state.model_name = config_settings['model_name']
if 'base_url' not in st.session_state:
    st.session_state.base_url = config_settings['base_url']
if 'system_prompt' not in st.session_state:
    st.session_state.system_prompt = config_settings['system_prompt']
if 'text_content' not in st.session_state:
    st.session_state.text_content = ""

def extract_text_from_file(uploaded_file):
    """从不同格式的文件中提取文本"""
    text = ""
    file_extension = uploaded_file.name.split('.')[-1].lower()
    
    if file_extension == 'txt':
        text = uploaded_file.read().decode('utf-8')
    elif file_extension == 'docx':
        doc = docx.Document(uploaded_file)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    elif file_extension == 'pdf':
        pdf_reader = PyPDF2.PdfReader(uploaded_file)
        text = '\n'.join([page.extract_text() for page in pdf_reader.pages])
    elif file_extension in ['html', 'htm']:
        soup = BeautifulSoup(uploaded_file, 'html.parser')
        text = soup.get_text(separator='\n')
    elif file_extension == 'md':
        text = uploaded_file.read().decode('utf-8')
    elif file_extension == 'pptx':
        prs = Presentation(uploaded_file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        text = '\n'.join(text_runs)
    
    return text

def process_text_with_gpt(text, api_key, model_name, base_url):
    """使用GPT处理文本，生成标题和正文格式"""
    try:
        client = openai.OpenAI(
            api_key=api_key,
            base_url=base_url
        )
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": st.session_state.system_prompt},
                {"role": "user", "content": text}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"GPT处理失败: {str(e)}")
        return text

async def generate_audio(text, voice_id, progress_placeholder):
    """生成音频文件"""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_file:
        progress_placeholder.text("生成音频中...")
        tts = edge_tts.Communicate(text, voice_id, rate='+25%')
        await tts.save(temp_file.name)
        progress_placeholder.text("音频生成完毕！")
        return temp_file.name

def extract_text_from_url(url):
    """从URL获取网页内容并提取文本"""
    try:
        # 发送HTTP请求
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()  # 检查请求是否成功
        
        # 使用BeautifulSoup解析HTML
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # 移除script、style和其他不需要的元素
        for element in soup(["script", "style", "meta", "link", "noscript", "header", "footer", "nav", "aside"]):
            element.decompose()
        
        # 处理标题
        if soup.title:
            title = soup.title.string.strip()
        else:
            title = ""
        
        # 处理正文内容
        content_parts = []
        processed_text = set()  # 用于去重
        
        # 添加标题
        if title:
            content_parts.append(title + "\n\n")
            processed_text.add(title)
        
        # 首先处理主要内容区域（如果存在）
        main_content = soup.find(['main', 'article']) or soup
        
        # 处理段落和其他文本元素
        for element in main_content.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p'], recursive=True):
            text = element.get_text(strip=True)
            if text and text not in processed_text:  # 只处理非空且未处理过的文本
                processed_text.add(text)
                # 对于标题元素添加额外的换行
                if element.name.startswith('h'):
                    content_parts.append("\n" + text + "\n")
                # 对于段落添加适当的间距
                else:
                    content_parts.append(text + "\n\n")
        
        # 处理可能的其他有意义的div内容
        for element in main_content.find_all('div', recursive=True):
            # 只处理直接包含文本的div，避免处理包含其他元素的div
            if element.find(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']) is None:
                text = element.get_text(strip=True)
                if text and len(text) > 50 and text not in processed_text:  # 只处理长度超过50的文本
                    processed_text.add(text)
                    content_parts.append(text + "\n\n")
        
        # 合并所有内容
        text = ''.join(content_parts)
        
        # 清理多余的空行和空格
        # 将连续的多个空行替换为两个空行
        text = re.sub(r'\n\s*\n', '\n\n', text)
        # 删除行首和行尾的空格
        text = '\n'.join(line.strip() for line in text.split('\n'))
        # 确保段落之间有适当的间距
        text = re.sub(r'\n\n\n+', '\n\n', text)
        
        return text.strip()
    except Exception as e:
        raise Exception(f"获取网页内容失败: {str(e)}")

# Streamlit UI
st.title("文本转语音生成器")

# GPT设置
st.sidebar.title("GPT设置")
use_gpt = st.sidebar.checkbox("使用GPT处理文本", value=st.session_state.use_gpt)
if use_gpt != st.session_state.use_gpt:
    st.session_state.use_gpt = use_gpt
    save_config({
        'use_gpt': use_gpt,
        'model_name': st.session_state.model_name,
        'base_url': st.session_state.base_url,
        'api_key': st.session_state.openai_api_key or "",
        'system_prompt': st.session_state.system_prompt
    })

if use_gpt:
    # Model选择
    st.sidebar.subheader("模型设置")
    models = ["google/gemini-2.0-flash-exp:free", "glm-4-flash", "deepseek-chat"]
    use_preset_model = st.sidebar.checkbox("使用预设模型", value=st.session_state.model_name in models)
    
    if use_preset_model:
        model_name = st.sidebar.selectbox(
            "选择预设模型",
            models,
            index=models.index(st.session_state.model_name) if st.session_state.model_name in models else 0
        )
    else:
        model_name = st.sidebar.text_input(
            "输入模型名称",
            value=st.session_state.model_name if st.session_state.model_name not in models else ""
        )
    
    if model_name != st.session_state.model_name:
        st.session_state.model_name = model_name
        save_config({
            'use_gpt': use_gpt,
            'model_name': model_name,
            'base_url': st.session_state.base_url,
            'api_key': st.session_state.openai_api_key or "",
            'system_prompt': st.session_state.system_prompt
        })
    
    # System Prompt输入
    system_prompt = st.sidebar.text_area(
        "System Prompt",
        value=st.session_state.system_prompt,
        height=100
    )
    if system_prompt != st.session_state.system_prompt:
        st.session_state.system_prompt = system_prompt
        save_config({
            'use_gpt': use_gpt,
            'model_name': model_name,
            'base_url': st.session_state.base_url,
            'api_key': st.session_state.openai_api_key or "",
            'system_prompt': system_prompt
        })
    
    # Base URL输入
    base_url = st.sidebar.text_input(
        "Base URL",
        value=st.session_state.base_url
    )
    if base_url != st.session_state.base_url:
        st.session_state.base_url = base_url
        save_config({
            'use_gpt': use_gpt,
            'model_name': model_name,
            'base_url': base_url,
            'api_key': st.session_state.openai_api_key or "",
            'system_prompt': st.session_state.system_prompt
        })
    
    # API Key输入
    api_key = st.sidebar.text_input("OpenAI API Key", 
                                   value=st.session_state.openai_api_key or "", 
                                   type="password")
    if api_key != st.session_state.openai_api_key:
        st.session_state.openai_api_key = api_key
        save_config({
            'use_gpt': use_gpt,
            'model_name': model_name,
            'base_url': base_url,
            'api_key': api_key,
            'system_prompt': st.session_state.system_prompt
        })

# URL输入
col1, col2 = st.columns([4, 1])
url = col1.text_input("输入网页URL")
if col2.button("获取内容"):
    if url:
        try:
            with st.spinner("正在获取网页内容..."):
                text_content = extract_text_from_url(url)
                st.session_state.text_content = text_content
                st.rerun()
        except Exception as e:
            st.error(str(e))
    else:
        st.warning("请输入URL")

# 文件上传
st.write("或者上传文件：")
uploaded_file = st.file_uploader("上传文件", type=['md', 'txt', 'html', 'htm', 'docx', 'pdf', 'pptx'])
text_content = st.session_state.text_content

if uploaded_file is not None:
    try:
        # 提取文本
        text_content = extract_text_from_file(uploaded_file)
        st.session_state.text_content = text_content
        
        # 使用GPT处理文本
        if use_gpt and st.session_state.openai_api_key:
            if st.button("处理文本"):
                with st.spinner("正在处理文本..."):
                    text_content = process_text_with_gpt(
                        text_content,
                        st.session_state.openai_api_key,
                        st.session_state.model_name,
                        st.session_state.base_url
                    )
                    st.session_state.text_content = text_content
        elif use_gpt and not st.session_state.openai_api_key:
            st.warning("请在侧边栏输入OpenAI API Key")
    except Exception as e:
        st.error(f"文件处理失败: {str(e)}")
        text_content = ""
        st.session_state.text_content = ""

# 文本编辑区
text_content = st.text_area("编辑文本", text_content, height=300)
st.session_state.text_content = text_content

# 整理文本按钮
col1, col2 = st.columns([1, 5])
if col1.button("整理文本"):
    if text_content:
        if use_gpt and st.session_state.openai_api_key:
            with st.spinner("正在整理文本..."):
                text_content = process_text_with_gpt(
                    text_content,
                    st.session_state.openai_api_key,
                    st.session_state.model_name,
                    st.session_state.base_url
                )
                # 更新文本编辑区
                st.session_state.text_content = text_content
                st.rerun()
        elif use_gpt and not st.session_state.openai_api_key:
            st.warning("请在侧边栏输入OpenAI API Key")
        else:
            st.warning("请在侧边栏启用GPT功能并完成相关配置")
    else:
        st.warning("请先输入文本内容")

# Voice selection
voices = [
    "zh-CN-YunjianNeural (Male)",
    "zh-CN-XiaoxiaoNeural (Female)", 
    "en-US-AvaMultilingualNeural (Female)", 
    "en-US-EmmaMultilingualNeural (Female)",
    "zh-CN-YunyangNeural (Male)", 
    "zh-CN-YunxiNeural (Male)",
    "en-US-BrianMultilingualNeural (Male)", 
    "en-US-RogerNeural (Male)",
    "en-US-SteffanNeural (Male)"
]

selected_voice = st.selectbox("选择语音", voices)
voice_id = selected_voice.split(" (")[0]

# 生成音频按钮
if st.button("生成音频"):
    if text_content:
        try:
            progress_placeholder = st.empty()
            
            # 生成音频
            audio_file = asyncio.run(generate_audio(text_content, voice_id, progress_placeholder))
            
            # 显示音频预览
            with open(audio_file, 'rb') as f:
                audio_bytes = f.read()
                st.audio(audio_bytes, format='audio/mp3')
                
                # 下载按钮
                st.download_button(
                    label="下载MP3",
                    data=audio_bytes,
                    file_name="generated_audio.mp3",
                    mime="audio/mpeg"
                )
            
            # 清理临时文件
            os.unlink(audio_file)
        except Exception as e:
            st.error(f"音频生成失败: {str(e)}")
    else:
        st.error("请先输入或上传文本内容") 